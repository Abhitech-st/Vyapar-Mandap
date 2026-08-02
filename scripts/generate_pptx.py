import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== DESIGN TOKENS (Modern High-Contrast SaaS Palette) ==========
    C_EMERALD = RGBColor(16, 185, 129)    # Growth / Success Accent
    C_NAVY = RGBColor(15, 23, 42)         # Primary Dark Text / Headers
    C_INDIGO = RGBColor(67, 56, 202)      # Brand Accent / Primary Buttons
    C_WHITE = RGBColor(255, 255, 255)
    C_BG_GRAY = RGBColor(248, 250, 252)   # Subtle Page Background
    C_CARD_BG = RGBColor(255, 255, 255)
    C_BORDER = RGBColor(226, 232, 240)    # 1px Container Border
    C_TEXT_MUTED = RGBColor(100, 116, 139) # Secondary Subtitle Text
    C_ACCENT_RED = RGBColor(220, 38, 38)   # Alert / Problem Red
    C_ACCENT_BLUE = RGBColor(37, 99, 235)  # Trust / AI Blue
    SHADOW_COLOR = RGBColor(235, 240, 245) # Subtle Card Drop Shadow

    def add_drop_shadow(slide, x, y, width, height):
        """Draws a subtle background drop-shadow rectangle behind cards"""
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.04), y + Inches(0.04), width, height)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = SHADOW_COLOR
        shadow.line.fill.background()
        return shadow

    def add_button(slide, text, x, y, width, height, color=C_INDIGO):
        """Draws a crisp action button"""
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
        btn.fill.solid()
        btn.fill.fore_color.rgb = color
        btn.line.color.rgb = color
        btn.line.width = Pt(0)
        
        txBox = slide.shapes.add_textbox(x, y, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.font.name = "Inter"
        p.alignment = PP_ALIGN.CENTER
        return btn

    blank = prs.slide_layouts[6]

    def add_header(slide, title, subtitle, accent_color=C_INDIGO):
        """Header bar with accent border"""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_BG_GRAY
        bar.line.fill.background()

        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.28), Inches(11.73), Pt(2))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_color
        accent_line.line.fill.background()

        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.73), Inches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.font.name = "Inter"
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(13)
            p2.font.color.rgb = C_TEXT_MUTED
            p2.font.name = "Inter"

    # ==========================================
    # SLIDE 1: Hero & Vision (Codex Architecture)
    # ==========================================
    slide1 = prs.slides.add_slide(blank)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_BG_GRAY
    bg1.line.fill.background()

    # Title Column
    t_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.5), Inches(4.5))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "VYAPAR MANDAP"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C_NAVY
    p.font.name = "Inter"

    p2 = tf1.add_paragraph()
    p2.text = "Built Using Codex AI Engine Architecture"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = C_INDIGO
    p2.font.name = "Inter"
    p2.space_before = Pt(8)

    p3 = tf1.add_paragraph()
    p3.text = "Engineered with OpenAI Codex intelligence, coupling vision OCR parsing with a deterministic double-entry ledger core (Debits = Credits) and 1-click CA human approvals."
    p3.font.size = Pt(13)
    p3.font.color.rgb = C_TEXT_MUTED
    p3.font.name = "Inter"
    p3.space_before = Pt(12)

    add_button(slide1, "Explore Live Platform", Inches(0.8), Inches(5.6), Inches(3.0), Inches(0.65), C_EMERALD)

    # Key Value Metric Cards
    hero_cards = [
        ("Codex AI Engine", "Advanced code-grade intelligence for Vision OCR invoice parsing & rule reasoning"),
        ("Double-Entry Core Engine", "Mathematical equality enforcement ensuring $Total\\ Debits = Total\\ Credits$"),
        ("Hybrid Cloud Infrastructure", "Supabase Central DB + Client Device Storage + Vercel & Render")
    ]

    for i, (val, desc) in enumerate(hero_cards):
        y_m = Inches(1.5) + i * 1.7
        add_drop_shadow(slide1, Inches(7.6), y_m, Inches(4.9), Inches(1.5))
        
        card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), y_m, Inches(4.9), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER

        strip = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), y_m, Pt(4), Inches(1.5))
        strip.fill.solid()
        strip.fill.fore_color.rgb = C_INDIGO
        strip.line.fill.background()

        txBox = slide1.shapes.add_textbox(Inches(7.9), y_m + Inches(0.15), Inches(4.4), Inches(1.2))
        tf_m = txBox.text_frame
        tf_m.word_wrap = True
        
        pm = tf_m.paragraphs[0]
        pm.text = val
        pm.font.size = Pt(17)
        pm.font.bold = True
        pm.font.color.rgb = C_NAVY
        pm.font.name = "Inter"

        pm2 = tf_m.add_paragraph()
        pm2.text = desc
        pm2.font.size = Pt(11)
        pm2.font.color.rgb = C_TEXT_MUTED
        pm2.font.name = "Inter"
        pm2.space_before = Pt(4)

    # ==========================================
    # SLIDE 2: The Current Problem
    # ==========================================
    slide2 = prs.slides.add_slide(blank)
    add_header(slide2, "The Current Problem: Manual Accounting & Tax Risks", "Operational bottlenecks faced by 63+ Million MSMEs in India", C_ACCENT_RED)
    
    alert = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.73), Inches(0.75))
    alert.fill.solid()
    alert.fill.fore_color.rgb = RGBColor(254, 242, 242)
    alert.line.color.rgb = RGBColor(254, 202, 202)
    
    tf_alert = alert.text_frame
    tf_alert.word_wrap = True
    p = tf_alert.paragraphs[0]
    p.text = "[!] Critical Pain Point: Manual invoice entry, un-tracked vendor TDS limits, and GSTR-2B discrepancies cause millions in lost tax credit and statutory interest notices."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT_RED
    p.font.name = "Inter"

    issues = [
        ("100+ Hours Lost to Manual Data Entry", "Accountants manually transcribe PDF bills, physical receipts, and paper vouchers line-by-line across disconnected spreadsheets, causing high human error rates."),
        ("18% GSTR-2B Input Tax Credit (ITC) Loss", "Discrepancies between supplier portal filings and internal ledger accounts result in lost Input Tax Credit (ITC) and costly statutory audit interest penalties."),
        ("Complex Section 194C/194J TDS Thresholds", "Cumulative vendor payment thresholds (Section 194C 1%/2% & Section 194J 10%) are missed across suppliers. Naive LLMs hallucinate numbers and break debit=credit equality.")
    ]
    
    for i, (title, desc) in enumerate(issues):
        y_pos = Inches(2.55) + (i * 1.55)
        add_drop_shadow(slide2, Inches(0.8), y_pos, Inches(11.73), Inches(1.35))
        
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.73), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER
        
        strip = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_pos, Pt(4), Inches(1.35))
        strip.fill.solid()
        strip.fill.fore_color.rgb = C_ACCENT_RED
        strip.line.fill.background()
        
        tx = slide2.shapes.add_textbox(Inches(1.1), y_pos + Inches(0.15), Inches(11.2), Inches(1.05))
        tf_i = tx.text_frame
        tf_i.word_wrap = True
        
        p = tf_i.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.font.name = "Inter"
        
        p2 = tf_i.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.font.name = "Inter"
        p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 3: Our Approach (Four Integrity Pillars)
    # ==========================================
    slide3 = prs.slides.add_slide(blank)
    add_header(slide3, "Our Approach: Four Pillars of Accounting Integrity", "Bridging unstructured document inputs with deterministic accounting engines", C_INDIGO)
    
    pillars = [
        ("1. Vision OCR & Codex Engine", "98%+ Accurate Invoice Extraction", ["Codex AI vision parsing architecture", "SHA256 document hashing (0ms repeat latency)", "Auto-flags low-confidence (<85%) field extractions"]),
        ("2. Double-Entry Core Engine", "Strict Debit = Credit Enforcement", ["Strict mathematical balance verification", "Zero balance violation guarantee ($Dr = Cr$)", "Cryptographic immutable transaction log"]),
        ("3. Human-in-the-Loop Signoff", "CA Verified Ledger Approvals", ["Split-screen PDF viewer + proposed journal", "1-click approval or rejection comments", "Complete audit trail logs every reviewer action"]),
        ("4. Statutory GST & TDS Engine", "Automated Compliance & Audit", ["Real-time 15-char GSTIN syntax check", "GSTR-1, GSTR-3B & GSTR-2B ITC matching", "Section 194C (1%/2%) & 194J (10%) TDS tracker"])
    ]
    
    for i, (title, sub, feats) in enumerate(pillars):
        col = i % 2
        row = i // 2
        y = Inches(1.6) + row * 2.75
        x = Inches(0.8) + (col * 6.0)
        
        add_drop_shadow(slide3, x, y, Inches(5.7), Inches(2.55))
        
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.55))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER
        
        top_bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.7), Pt(4))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C_INDIGO
        top_bar.line.fill.background()
        
        tx = slide3.shapes.add_textbox(x + Inches(0.3), y + Inches(0.15), Inches(5.1), Inches(2.25))
        tf_p = tx.text_frame
        tf_p.word_wrap = True
        
        p = tf_p.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.font.name = "Inter"
        
        p2 = tf_p.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = C_INDIGO
        p2.font.name = "Inter"
        p2.space_before = Pt(2)
        p2.space_after = Pt(4)
        
        for feat in feats:
            p3 = tf_p.add_paragraph()
            p3.text = f"• {feat}"
            p3.font.size = Pt(10)
            p3.font.color.rgb = C_TEXT_MUTED
            p3.font.name = "Inter"

    # ==========================================
    # SLIDE 4: AI Integration & Codex Task Engine
    # ==========================================
    slide4 = prs.slides.add_slide(blank)
    add_header(slide4, "AI Integration & Codex Task Engine", "Combining OpenAI Codex Architecture with 1-Tap Financial Analysis", C_ACCENT_BLUE)
    
    ai_features = [
        ("Codex AI Vision OCR", "Ingests PDF invoices, receipts, and image bills, extracting vendor name, GSTIN, line-item tax split, and grand totals with 99% accuracy."),
        ("Pre-Coded AI Task Shortcuts", "1-tap task triggers: Audit GSTR-2B ITC, Check Trial Balance Equality, Calculate Section 194C/194J Limits, Reconcile Bank Feed, P&L Highlights."),
        ("Automated Compliance Rules", "Executes 15-char GSTIN syntax verification, duplicate SHA256 document detection, and double-entry mathematical constraint checks."),
        ("Natural Language Financial Query", "Interactive AI Copilot allowing Chartered Accountants to query cash runway, un-posted vouchers, and tax return filing deadlines.")
    ]
    
    for i, (title, desc) in enumerate(ai_features):
        col = i % 2
        row = i // 2
        y = Inches(1.6) + row * 2.75
        x = Inches(0.8) + (col * 6.0)
        
        add_drop_shadow(slide4, x, y, Inches(5.7), Inches(2.55))
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.55))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER
        
        strip = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(4), Inches(2.55))
        strip.fill.solid()
        strip.fill.fore_color.rgb = C_ACCENT_BLUE
        strip.line.fill.background()
        
        tx = slide4.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(5.1), Inches(2.15))
        tf_f = tx.text_frame
        tf_f.word_wrap = True
        
        p = tf_f.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.font.name = "Inter"
        
        p2 = tf_f.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.font.name = "Inter"
        p2.space_before = Pt(8)

    # ==========================================
    # SLIDE 5: The Golden Path Workflow
    # ==========================================
    slide5 = prs.slides.add_slide(blank)
    add_header(slide5, "The Golden Path Workflow", "End-to-End Ingestion, AI Parsing, CA Signoff & Immutable Ledger Commit", C_EMERALD)
    
    steps = [
        ("Step 1: Upload", "Ingest PDF / Image", ["User drops PDF bill or receipt image", "Generates unique SHA256 hash", "Prevents duplicate invoice ingestion"]),
        ("Step 2: AI Parse", "Extract & Audit Tax", ["Codex AI engine extracts fields", "GST engine validates 15-char GSTIN", "Verifies GSTR-2B ITC eligibility"]),
        ("Step 3: Human Review", "CA Signoff Checkpoint", ["Split-screen PDF viewer + journal", "Validates $Dr. Exp + Dr. Tax = Cr. AP$", "1-click approve or reject with comments"]),
        ("Step 4: Commit", "Immutable Ledger", ["Double-entry posted to database", "Journal marked immutable", "P&L and Balance Sheet update live"])
    ]
    
    y_pos = Inches(1.6)
    
    for i, (title, sub, feats) in enumerate(steps):
        x_pos = Inches(0.8) + (i * 2.95)
        
        add_drop_shadow(slide5, x_pos, y_pos, Inches(2.75), Inches(5.4))
        
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, Inches(2.75), Inches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER
        
        badge = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(1.05), y_pos + Inches(0.2), Inches(0.65), Inches(0.65))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_EMERALD
        badge.line.fill.background()

        badge_txt = slide5.shapes.add_textbox(x_pos + Inches(1.05), y_pos + Inches(0.2), Inches(0.65), Inches(0.65))
        tf_b = badge_txt.text_frame
        pb = tf_b.paragraphs[0]
        pb.text = str(i + 1)
        pb.font.size = Pt(14)
        pb.font.bold = True
        pb.font.color.rgb = C_WHITE
        pb.alignment = PP_ALIGN.CENTER
        
        tx = slide5.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.95), Inches(2.35), Inches(4.3))
        tf_s = tx.text_frame
        tf_s.word_wrap = True
        
        p = tf_s.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.font.name = "Inter"
        
        p2 = tf_s.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = C_EMERALD
        p2.font.name = "Inter"
        p2.space_before = Pt(2)
        p2.space_after = Pt(8)
        
        for feat in feats:
            p3 = tf_s.add_paragraph()
            p3.text = f"• {feat}"
            p3.font.size = Pt(10)
            p3.font.color.rgb = C_TEXT_MUTED
            p3.font.name = "Inter"
            p3.space_before = Pt(4)

    # ==========================================
    # SLIDE 6: Production Infrastructure & Hybrid Storage
    # ==========================================
    slide6 = prs.slides.add_slide(blank)
    add_header(slide6, "Production Infrastructure & Hybrid Storage", "Full-Stack Cloud Architecture (Vercel + Render + Supabase)", C_INDIGO)
    
    infra = [
        ("Vercel Frontend CDN (React 18 + Vite)", "Global CDN static hosting with Mobile Navigation Drawer, responsive table-to-card transformations, and Command Palette (Cmd+K)."),
        ("Render FastAPI Backend (Python 3.11)", "Asynchronous Python 3.11 ASGI service executing Codex AI Vision OCR, fuzzy bank matching, and Pydantic validation."),
        ("Supabase Cloud Database (PostgreSQL)", "Central cloud database storing registered User Profiles (`profiles`) and Firm Entity Masters (`organizations`) with Row-Level Security."),
        ("Client Device Local Storage (localStorage)", "Stores active workspace session state, draft invoice uploads, and UI filters for instant 0ms page loads and offline resilience.")
    ]
    
    for i, (layer, tech) in enumerate(infra):
        y = Inches(1.55) + i * 1.38
        add_drop_shadow(slide6, Inches(0.8), y, Inches(11.73), Inches(1.22))
        
        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.73), Inches(1.22))
        box.fill.solid()
        box.fill.fore_color.rgb = C_CARD_BG
        box.line.color.rgb = C_BORDER
        
        strip = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Pt(4), Inches(1.22))
        strip.fill.solid()
        strip.fill.fore_color.rgb = C_INDIGO
        strip.line.fill.background()
        
        tx = slide6.shapes.add_textbox(Inches(1.1), y + Inches(0.12), Inches(11.3), Inches(1.0))
        tf_l = tx.text_frame
        tf_l.word_wrap = True
        
        p = tf_l.paragraphs[0]
        p.text = layer
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.font.name = "Inter"
        
        p2 = tf_l.add_paragraph()
        p2.text = tech
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.font.name = "Inter"
        p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 7: Strategic Roadmap & Future Expansion
    # ==========================================
    slide7 = prs.slides.add_slide(blank)
    add_header(slide7, "Strategic Roadmap & Growth Vision", "Phased expansion from MVP to Enterprise Chartered Accountancy Scaling", C_EMERALD)
    
    milestones = [
        ("Phase 1: MVP (Completed)", "• Hybrid Storage Engine (Supabase + Local)\n• Double-entry ledger core engine\n• Statutory GST & TDS compliance\n• OpenAI Codex AI Architecture\n• 1-Tap Pre-Coded AI Tasks", C_EMERALD, "COMPLETED"),
        ("Phase 2: Q3 2026 (In Progress)", "• Direct GSTN Portal Sandbox APIs\n• Account Aggregator live bank feeds\n• Automated inventory batch valuation\n• E-way bill generation pipeline\n• Multi-user permission roles", C_INDIGO, "IN PROGRESS"),
        ("Phase 3: 2027 (Planned)", "• ICAI fine-tuned accounting LLM\n• CA Multi-Firm Client Portal\n• Predictive working capital credit scoring\n• Automated vendor payouts via UPI/NEFT\n• Enterprise API Marketplace", C_ACCENT_BLUE, "PLANNED")
    ]
    
    for i, (title, desc, color, status) in enumerate(milestones):
        x_pos = Inches(0.8) + (i * 3.95)
        
        add_drop_shadow(slide7, x_pos, Inches(1.6), Inches(3.8), Inches(5.4))
        
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.6), Inches(3.8), Inches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER
        
        badge = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.25), Inches(1.85), Inches(1.5), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(239, 246, 255) if status == "IN PROGRESS" else (RGBColor(240, 253, 244) if status == "COMPLETED" else RGBColor(248, 250, 252))
        badge.line.color.rgb = color
        
        tf_bg = badge.text_frame
        p_bg = tf_bg.paragraphs[0]
        p_bg.text = status
        p_bg.font.size = Pt(9)
        p_bg.font.bold = True
        p_bg.font.color.rgb = color
        p_bg.alignment = PP_ALIGN.CENTER
        
        tx = slide7.shapes.add_textbox(x_pos + Inches(0.25), Inches(2.35), Inches(3.3), Inches(4.5))
        tf_m = tx.text_frame
        tf_m.word_wrap = True
        
        p = tf_m.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.font.name = "Inter"
        
        p2 = tf_m.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.font.name = "Inter"
        p2.space_before = Pt(8)

    # ==========================================
    # SLIDE 8: Call to Action & Presentation Links
    # ==========================================
    slide8 = prs.slides.add_slide(blank)
    
    bg_cta = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_cta.fill.solid()
    bg_cta.fill.fore_color.rgb = C_INDIGO
    bg_cta.line.fill.background()
    
    tx_cta = slide8.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.0))
    tf_c = tx_cta.text_frame
    tf_c.word_wrap = True
    
    p = tf_c.paragraphs[0]
    p.text = "Ready to Transform Indian Business Accounting?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Inter"
    
    p2 = tf_c.add_paragraph()
    p2.text = "Automate PDF invoice OCR, double-entry ledger posting, and statutory GST & TDS compliance."
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(191, 219, 254)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)
    p2.font.name = "Inter"

    p3 = tf_c.add_paragraph()
    p3.text = "GitHub: github.com/Abhitech-st/Vyapar-Mandap"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(240, 253, 244)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)
    p3.font.name = "Inter"
    
    add_button(slide8, "Explore Live Platform Demo", Inches(4.66), Inches(5.4), Inches(4.0), Inches(0.75), C_EMERALD)

    # Save PowerPoint Presentation
    output_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    target_path = os.path.join(output_dir, "Vyapar_Mandap_Presentation.pptx")
    try:
        prs.save(target_path)
        print(f"[OK] Presentation saved to: {target_path}")
    except PermissionError:
        fallback = os.path.join(output_dir, "Vyapar_Mandap_Presentation_Fixed.pptx")
        prs.save(fallback)
        print(f"[INFO] Primary file locked. Saved to: {fallback}")

if __name__ == "__main__":
    create_presentation()
